import streamlit as st
import pandas as pd
import geopandas as gpd
import folium
from folium.features import GeoJsonTooltip
import branca
from folium.plugins import TagFilterButton
from streamlit_folium import folium_static
from folium.plugins import Fullscreen
from collections import Counter
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import base64
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode, JsCode
from streamlit_option_menu import option_menu
from io import BytesIO
import PyPDF2
from openai import OpenAI
# from openai.types.chat import ChatCompletionMessage
# from openai.types.chat.chat_completion import ChatCompletion, Choice
# import tiktoken
from pptx import Presentation
from docx import Document
from streamlit_chat import message
from thefuzz import process

levl0=gpd.read_file("europe.geojson")
levl2=gpd.read_file("NUTS_2_Q2.geojson")
levl1=gpd.read_file("NUTS_1_Q1.geojson")
levl3=gpd.read_file("NUTS_3_Q1.geojson")
levl3.drop_duplicates(subset="NUTS_ID", keep="first", inplace=True)
zip_code=pd.read_excel("zipcode_nuts with uk instead of gb.xlsx")
dsv=pd.read_excel("DSV Branches.xlsx")

st.set_page_config(layout='wide')

# Function to load data

def load_data():
    
    uploaded_file = st.session_state.get('uploaded_file', None)
    
    if uploaded_file is not None:
        data = pd.read_excel(uploaded_file)
        if 'Date' not in data.columns:
            data['Date'] = pd.Series(dtype='datetime64[ns]')  
        if "Category1" in data.columns and "Category2" in data.columns:
            data=data[["ZC from","ZC to","Date","kg","ldm","m3","Branch",'Category1','Category2']]
        else:
            data=data[["ZC from","ZC to","kg","ldm","m3","Branch","Date"]]

        data['Branch'] = data['Branch'].fillna('undefined')
        

        data[['kg', 'ldm', 'm3']] = data[['kg', 'ldm', 'm3']].fillna(0)
        data=data[~(data["ZC from"].isna() | data["ZC to"].isna() | ( (( (data["kg"] == 0)) & ( (data["ldm"] == 0)) & ( (data["m3"] == 0)))))]
        
        country_codes = ['AD', 'ME', 'RU', 'XK', 'UA', 'BY', 'BA']
        for code in country_codes:
            data.loc[data['ZC from'].str[:2] == code, 'ZC from'] = code
            data.loc[data['ZC to'].str[:2] == code, 'ZC to'] = code
        data['ZC from'] = data['ZC from'].apply(lambda x: 'UK' + x[2:] if x.startswith('GB') else x)
        data['ZC to'] = data['ZC to'].apply(lambda x: 'UK' + x[2:] if x.startswith('GB') else x)
        
        
        if not pd.api.types.is_datetime64_any_dtype(data['Date']):
            data['Date'] = pd.to_datetime(data['Date'], errors='coerce') 
        data['Date'] = data['Date'].dt.date
       

            
        
        return data
    else:
        st.error("please upload your data first",icon="🚨")
def process_data(data, selected_dsv_country, selected_parcel, selected_grp, selected_ltl, pw_ldm, pw_cbm,factor,input_factor):
      
             
    
    data['Cntry from'] = data['ZC from'].str[:2]
    data['Cntry to'] = data['ZC to'].str[:2]

    
    data['Way'] = data.apply(lambda row: 'Exp' if row['Cntry from'] == selected_dsv_country and row['Cntry to'] != selected_dsv_country
                            else 'Imp' if row['Cntry to'] == selected_dsv_country and row['Cntry from'] != selected_dsv_country
                            else 'Dom' if row['Cntry from'] == row['Cntry to'] == selected_dsv_country
                            else 'X-trade', axis=1)

    
    data['kg'] = pd.to_numeric(data['kg'], errors='coerce')


    
    data['ldm'] = pd.to_numeric(data['ldm'], errors='coerce')
    data['m3'] = pd.to_numeric(data['m3'], errors='coerce')
    if factor=="Yes":
        
            data["ldm_eq"]=data.apply(lambda row: max(row['ldm'], row['kg'] / float(input_factor)), axis=1)
            data["ldm"]=data["ldm_eq"]
            st.session_state.factor_phrase= (f" |kg/ldm factor ={st.session_state.input_factor}")
    else:
        data=data
        st.session_state.factor_phrase= " "
    
    data['PW DSV'] = data.apply(lambda row: max(row['kg'], row['ldm'] * float(pw_ldm), row['m3'] * float(pw_cbm)), axis=1)

    
    data['Bracket'] = data['PW DSV'].apply(
        lambda x: 30 if 0 < x <= 30
        else 100 if 30 < x <= 100
        else 250 if x <= 250
        else 500 if x <= 500
        else 1000 if x <= 1000
        else 2500 if x <= 2500
        else 5000 if x <= 5000
        else 7500 if x <= 7500
        else 10000 if x <= 10000
        else 15000 if x <= 15000
        else 20000 if x <= 20000
        else 25500
    )
    data['Product'] = data['PW DSV'].apply(
    lambda x: 'Parcel' if x <= float(selected_parcel)
    else 'GRP' if x <= float(selected_grp)
    else 'LTL' if x <= float(selected_ltl)
    else 'FTL'
    )
    


    return data
    
dsv_country=["AL","AT","BA","BE","BG","CH","CZ","DE","DK","EE","ES","FI","FR","GB","GR","HR","HU","IE","IT","LT","LU","LV","ME","MK","NL","NO","PL","PT","RO","RS","SE","SI","SK","TR","NI","XK"]
brakets=["30","100","250","500","1000","2500","5000","7500","10000","15000","20000","25000"]
if 'selected' not in st.session_state:
    st.session_state.selected = "Upload data"

selected_option  = option_menu(
menu_title=None,
options=["Upload data", "Shipment Summary", "Shipment Profile","Maps","Collection Analysis","Regularity Detector","Document","Data cleaning"],
icons=["bi-cloud-upload", "bi bi-bar-chart-fill", "graph-up","bi bi-globe-europe-africa","bi bi-calendar-event","bi bi-filter","bi bi-chat-dots"],
menu_icon="cast",
default_index=0,
orientation="horizontal",
)

if selected_option  != st.session_state.selected:
     st.session_state.selected = selected_option 

# Page for uploading and viewing the Excel file
if st.session_state.selected == "Upload data":
    st.title('Upload your data')
    uploaded_file = st.file_uploader("Put your data here", type=['xlsx'])
    
    col_empty=st.empty()
    form_empty=st.empty()
    with col_empty:
        col1,col2=st.columns([1.5,3],gap='large')
        with col1:
            
                st.write("Here's an example of how you data should look like:")
                st.image("Capture d’écran 2024-08-08 093822.png")
                st.write("*Column 'Date' is optional.")
                st.write("If needed you can also add columns 'Category1' and 'Category2' :")
                st.image("Capture d’écran 2024-09-16 170227.png")
        with col2:
            if uploaded_file is not None:
                st.session_state['uploaded_file'] = uploaded_file
                
                data = pd.read_excel(uploaded_file)
                
                
                if 'Date' not in data.columns:
                    data['Date'] = pd.Series(dtype='datetime64[ns]')  
                if "Category1" in data.columns and "Category2" in data.columns:
                    data=data[["ZC from","ZC to","Date","kg","ldm","m3","Branch",'Category1','Category2']]
                else:
                    data=data[["ZC from","ZC to","kg","ldm","m3","Branch","Date"]]
                st.write("")
                st.dataframe(data,height=250, use_container_width=True)
    
    if uploaded_file is not None:        
        
    
        with form_empty:
            with st.form("filters"):
                st.write("**Basic options**")
                col3,col35,col4,col5,col6=st.columns([2,1.5,2,2,2],gap="medium")
                with col3:
                    st.write("Some information depends on country perspective, select your country:")
                    st.session_state.selected_dsv_country=st.selectbox("",dsv_country,index=dsv_country.index("FR"))
                with col35:
                    st.write("Choose how payweight should be calculated:")
                    col01,col02=st.columns([1,1])
                    with col01:
                        st.session_state.pw_ldm=st.number_input("kg/ldm :",value=1750)
                    with col02:
                        st.session_state.pw_cbm=st.number_input("kg/cbm:",value=330)
                with col4:
                    st.write("Define parcel's weight(Kg), up to:")
                    st.session_state.selected_parcel=st.select_slider("",brakets,value="30")
                with col5:
                    st.write("Define groupage's weight(Kg), up to:")
                    st.session_state.selected_grp=st.select_slider("",brakets,value="2500")
                with col6:
                    st.write("Define LTL's weight(Kg), up to:")
                    st.session_state.selected_ltl=st.select_slider("",brakets,value="20000")
                
                st.write("")
                st.write("**Advanced option(use only if you know what you are doing)**")    
                st.session_state.factor=st.radio("In case of  missing ldm or cbm, you can apply a kg/ldm factor. Do you want to apply one ?",["No","Yes"])
                col8,col9=st.columns([1,8])
                with col8:
                    st.session_state.input_factor = st.number_input("If yes, enter your kg/ldm factor:",value=1750)
                col11,col12,col13=st.columns([1.5,1,1])
                with col12:
                        st.write("")
                        submitted = st.form_submit_button("Apply modifications")
        
                    
            if submitted:
                col_empty.empty()
                form_empty.empty()
                with col_empty:
                    processed_data =process_data(load_data(), st.session_state.selected_dsv_country, st.session_state.selected_parcel, st.session_state.selected_grp, st.session_state.selected_ltl, st.session_state.pw_ldm, st.session_state.pw_cbm,st.session_state.factor,st.session_state.input_factor)
                    processed_data=processed_data[~(processed_data["ZC from"].isna() | processed_data["ZC to"].isna() | ( ((processed_data["kg"].isna()| (processed_data["kg"] == 0)) & (processed_data["ldm"].isna()| (processed_data["ldm"] == 0)) & (processed_data["m3"].isna()| (processed_data["m3"] == 0)))))]
                    st.dataframe(processed_data ,use_container_width=True)
                    st.session_state.processed_data = processed_data
        
                    
            else:
                st.session_state.processed_data =process_data(load_data(), st.session_state.selected_dsv_country, st.session_state.selected_parcel, st.session_state.selected_grp, st.session_state.selected_ltl, st.session_state.pw_ldm, st.session_state.pw_cbm,st.session_state.factor,st.session_state.input_factor)
                st.session_state.processed_data =st.session_state.processed_data[~(st.session_state.processed_data["ZC from"].isna() | st.session_state.processed_data["ZC to"].isna() | ( ((st.session_state.processed_data["kg"].isna()| (st.session_state.processed_data["kg"] == 0)) & (st.session_state.processed_data["ldm"].isna()| (st.session_state.processed_data["ldm"] == 0)) & (st.session_state.processed_data["m3"].isna()| (st.session_state.processed_data["m3"] == 0)))))]
        removed_rows=data[(data["ZC from"].isna() | data["ZC to"].isna() |( ( ((data["kg"] == 0)|data["kg"].isna()) &  ((data["ldm"] == 0)|data["ldm"].isna()) &  ((data["m3"] == 0)|data["m3"].isna()))))]
        if not removed_rows.empty:
            lengh=len(removed_rows)
            st.write(f"{lengh} rows have been removed due to missing data:")
            st.dataframe(removed_rows, use_container_width=True)
            def to_excel(df):
                output = BytesIO()
                with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                    df.to_excel(writer, index=False, sheet_name='Sheet1')
                return output.getvalue()

            # Convert the cleaned DataFrame to an Excel file
            excel_data = to_excel(removed_rows)

            # Download button
            st.download_button(
                label="Download missing rows as Excel",
                data=excel_data,
                file_name='missing_data.xlsx',
                mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )
def filters_with_categories(data):
                        st.write("**Filters option**")
                        selected_branch = st.multiselect('Select branch', options=data['Branch'].unique())
                        if selected_branch:
                            data = data[data['Branch'].isin(selected_branch)]
                        else:
                            data = data
                        selected_category1=st.multiselect('Select Category 1', options=data['Category1'].unique())
                        if selected_category1:
                            data = data[data['Category1'].isin(selected_category1)]
                        else:
                            data=data

                        selected_category2=st.multiselect('Select Category 2', options=data['Category2'].unique())
                        if selected_category2:
                            data = data[data['Category2'].isin(selected_category2)]
                        else:
                            data=data
                    
                        filtered_cntry_from = sorted(data['Cntry from'].unique())
                        selected_cntry_from = st.multiselect('Select Country From', filtered_cntry_from)
                        if selected_cntry_from:
                            data = data[data['Cntry from'].isin(selected_cntry_from)]
                        else:
                            data = data  

                        
                        filtered_zc_from = sorted(data['ZC from'].unique())
                        selected_zc_from = st.multiselect('Select Zip Code From', filtered_zc_from)

                        filtered_cntry_to = sorted(data['Cntry to'].unique())
                        selected_cntry_to = st.multiselect('Select Country To', filtered_cntry_to)
                        if selected_cntry_to:
                            data = data[data['Cntry to'].isin(selected_cntry_to)]
                        else:
                            data = data  

                    
                        filtered_zc_to = sorted(data['ZC to'].unique())
                        selected_zc_to = st.multiselect('Select Zip Code To', filtered_zc_to)
                        selected_product = st.multiselect('Select type of product', options=data['Product'].unique())
                        selected_way = st.multiselect('Select way', options=data['Way'].unique())

                        if data['Date'].isna().sum() == 0:
                            data['Date'] = pd.to_datetime(data['Date'])

                            first_date = data['Date'].min()
                            last_date = data['Date'].max()

                            start_date = st.date_input(
                                "Select start date",
                                value=first_date,
                                min_value=first_date,
                                max_value=last_date,
                                format="DD.MM.YYYY"
                                
                            )

                            end_date = st.date_input(
                                "Select end date",
                                value=last_date,
                                min_value=start_date,
                                max_value=last_date,
                                format="DD.MM.YYYY"
                            )
                        
                        
                            

                        if 'start_date' in locals() and 'end_date' in locals():
                            data = data[
                            (data['Cntry from'].isin(selected_cntry_from) if selected_cntry_from else data['Cntry from'].notnull()) &
                            (data['ZC from'].isin(selected_zc_from) if selected_zc_from else data['ZC from'].notnull()) &
                            (data['Cntry to'].isin(selected_cntry_to) if selected_cntry_to else data['Cntry to'].notnull()) &
                            (data['ZC to'].isin(selected_zc_to) if selected_zc_to else data['ZC to'].notnull())&
                            (data['Product'].isin(selected_product) if selected_product else data['Product'].notnull())&
                            (data['Way'].isin(selected_way) if selected_way else data['Way'].notnull())&
                            (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())&
                            (data['Date'] >= pd.to_datetime(start_date)) & (data['Date'] <= pd.to_datetime(end_date))&
                            (data['Category1'].isin(selected_category1) if selected_category1 else data['Category1'].notnull())&
                            (data['Category2'].isin(selected_category2) if selected_category2 else data['Category2'].notnull())]
                        else:
                            data = data[
                            (data['Cntry from'].isin(selected_cntry_from) if selected_cntry_from else data['Cntry from'].notnull()) &
                            (data['ZC from'].isin(selected_zc_from) if selected_zc_from else data['ZC from'].notnull()) &
                            (data['Cntry to'].isin(selected_cntry_to) if selected_cntry_to else data['Cntry to'].notnull()) &
                            (data['ZC to'].isin(selected_zc_to) if selected_zc_to else data['ZC to'].notnull())&
                            (data['Product'].isin(selected_product) if selected_product else data['Product'].notnull())&
                            (data['Way'].isin(selected_way) if selected_way else data['Way'].notnull())&
                            (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())&
                            (data['Category1'].isin(selected_category1) if selected_category1 else data['Category1'].notnull())&
                            (data['Category2'].isin(selected_category2) if selected_category2 else data['Category2'].notnull())]

                        return data
                
def apply_filters(data):
                        st.write("**Filters option**")
                        selected_branch = st.multiselect('Select branch', options=data['Branch'].unique())
                        if selected_branch:
                            data = data[data['Branch'].isin(selected_branch)]
                        else:
                            data = data
                    
                        filtered_cntry_from = sorted(data['Cntry from'].unique())
                        selected_cntry_from = st.multiselect('Select Country From', filtered_cntry_from)
                        if selected_cntry_from:
                            data = data[data['Cntry from'].isin(selected_cntry_from)]
                        else:
                            data = data  

                        
                        filtered_zc_from = sorted(data['ZC from'].unique())
                        selected_zc_from = st.multiselect('Select Zip Code From', filtered_zc_from)

                        filtered_cntry_to = sorted(data['Cntry to'].unique())
                        selected_cntry_to = st.multiselect('Select Country To', filtered_cntry_to)
                        if selected_cntry_to:
                            data = data[data['Cntry to'].isin(selected_cntry_to)]
                        else:
                            data = data  

                    
                        filtered_zc_to = sorted(data['ZC to'].unique())
                        selected_zc_to = st.multiselect('Select Zip Code To', filtered_zc_to)
                        selected_product = st.multiselect('Select type of product', options=data['Product'].unique())
                        selected_way = st.multiselect('Select way', options=data['Way'].unique())

                        if data['Date'].isna().sum() == 0:
                            data['Date'] = pd.to_datetime(data['Date'])

                            first_date = data['Date'].min()
                            last_date = data['Date'].max()

                            start_date = st.date_input(
                                "Select start date",
                                value=first_date,
                                min_value=first_date,
                                max_value=last_date,
                                format="DD.MM.YYYY"
                                
                            )

                            end_date = st.date_input(
                                "Select end date",
                                value=last_date,
                                min_value=start_date,
                                max_value=last_date,
                                format="DD.MM.YYYY"
                            )
                        
                        
                            

                        if 'start_date' in locals() and 'end_date' in locals():
                            data = data[
                            (data['Cntry from'].isin(selected_cntry_from) if selected_cntry_from else data['Cntry from'].notnull()) &
                            (data['ZC from'].isin(selected_zc_from) if selected_zc_from else data['ZC from'].notnull()) &
                            (data['Cntry to'].isin(selected_cntry_to) if selected_cntry_to else data['Cntry to'].notnull()) &
                            (data['ZC to'].isin(selected_zc_to) if selected_zc_to else data['ZC to'].notnull())&
                            (data['Product'].isin(selected_product) if selected_product else data['Product'].notnull())&
                            (data['Way'].isin(selected_way) if selected_way else data['Way'].notnull())&
                            (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())&
                            (data['Date'] >= pd.to_datetime(start_date)) & (data['Date'] <= pd.to_datetime(end_date))]
                        else:
                            data = data[
                            (data['Cntry from'].isin(selected_cntry_from) if selected_cntry_from else data['Cntry from'].notnull()) &
                            (data['ZC from'].isin(selected_zc_from) if selected_zc_from else data['ZC from'].notnull()) &
                            (data['Cntry to'].isin(selected_cntry_to) if selected_cntry_to else data['Cntry to'].notnull()) &
                            (data['ZC to'].isin(selected_zc_to) if selected_zc_to else data['ZC to'].notnull())&
                            (data['Product'].isin(selected_product) if selected_product else data['Product'].notnull())&
                            (data['Way'].isin(selected_way) if selected_way else data['Way'].notnull())&
                            (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())]
                        return data

if st.session_state.selected == "Shipment Summary":
        
        data=load_data()
        data = st.session_state.processed_data
        data["PW DSV"] = pd.to_numeric(data["PW DSV"], errors='coerce')
        data= data[data["PW DSV"].notna()]

        
        col1,col2=st.columns([1,7],gap="large")         
        with col1:

                if "Category1" in data.columns and "Category2" in data.columns:
                    data=filters_with_categories(data)
                    
                else:     
                  data=apply_filters(data)  
        with col2:
            st.header("Shipment summary")
            st.write(f"***Point of view from {st.session_state.selected_dsv_country} |  Parcel ≤ {st.session_state.selected_parcel} kg | GRP ≤ {st.session_state.selected_grp} kg | LTL ≤ {st.session_state.selected_ltl} kg | FTL > {st.session_state.selected_ltl} kg | Ratios {st.session_state.pw_cbm}/cbm & {st.session_state.pw_ldm}/ldm  {st.session_state.factor_phrase}.***")
            dom=data["ZC from"][data["Way"]=="Dom"].count()
            exp=data["ZC from"][data["Way"]=="Exp"].count()
            imp=data["ZC from"][data["Way"]=="Imp"].count()
            X_trade=data["ZC from"][data["Way"]=="X-trade"].count()

            values_way=[dom,exp,imp,X_trade]
            labels_way=["Dom","Exp","Imp","X_trade"]

            col1, col2,col3,col4 = st.columns([1.3,1.3,2,2])
            with col1:
                produit=data["Product"].unique().tolist()
                labels = []
                sh_values = []
                dict_product={}
                for k in produit:
                    summ=(data["Product"]==k).sum()
                    dict_product[f'nbr_{k}']=summ
                    sh_values.append(summ)
                    labels.append(k)
                colors =['#002664','#5D7AB5','#A9BCE2']
                fig = make_subplots(rows=1, specs=[[{'type':'domain'}]])
                fig.add_trace(go.Pie(labels=labels, values=sh_values, name="nbr of shipment"),1,1)
                fig.update_traces(hole=.5,marker=dict(colors=colors))
                fig.update_layout(
                annotations=[dict(text='Shipments', x=0.27, y=0.5, font_size=20, showarrow=False)])
                st.write("<h5><b>Product</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig,use_container_width=True)
                
            with col2:
                fig = make_subplots(rows=1, cols=1, specs=[[{'type':'domain'}]])
                fig.add_trace(go.Pie(labels=labels_way, values=values_way, name="Way"),1,1)
                fig.update_traces(hole=.5,marker=dict(colors=['#002664','#5D7AB5','#A9BCE2','#000000']))
                fig.update_layout(annotations=[dict(text='Way', x=0.5, y=0.5, font_size=20, showarrow=False)])
                st.write("<h5><b>Type</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig,use_container_width=True)
                
            
            with col3:
                
                my_list = sorted(data["Bracket"].tolist())
                count = Counter(my_list)
                total_items = sum(count.values())
                count_percentage_list = [(item, count, f"{round((count / total_items) * 100, 2)}%") for item, count in count.items()]
                df_bracket = pd.DataFrame(count_percentage_list, columns=['bracket', 'Count', 'percentage'])
                fig = px.bar(df_bracket, x='bracket', y='Count',color='Count', 
                color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'], text='percentage')
                fig.update_layout(
                    # title="Shipments per brackets ",
                    xaxis_title='',
                    yaxis_title='Shipments',
                    xaxis=dict(type='category'))
                fig.update_coloraxes(showscale=False)
                st.write("<h5><b>Shipments per brackets</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig,use_container_width=True)
                
                
            with col4:
                
                data["PW DSV"]=data["PW DSV"].astype(int)

                # creating the folium map
                data1=data.groupby(["ZC to"],as_index=False)["PW DSV"].count()
                data1=pd.merge(data1,zip_code,on='ZC to',how="left")
                data1['count'] = data1.groupby('ZC to')['ZC to'].transform('count')
                data1["PW DSV"]=(data1["PW DSV"])/data1["count"]
                data1=data1.groupby(["nuts0"],as_index=False)["PW DSV"].sum()
                merge=pd.merge(levl0,data1,right_on="nuts0" ,left_on="ISO2",how="right")
                
                m = folium.Map(location=[55.6761, 12.5683], zoom_start=2.5, zoom_control=False, tiles = "CartoDB Positron" )
                colums=["nuts0","PW DSV"]
                key="properties.ISO2"
                
                st.write("<h5><b>Countries</b></h5>", unsafe_allow_html=True)


                choropleth=folium.Choropleth(
                    geo_data=merge,
                    name="choropleth",
                    data=merge,
                    columns=colums,
                    key_on=key,
                    fill_color='Blues',
                    fill_opacity=0.7,
                    legend=False,
                    highlight=True,
                    ).geojson.add_to(m)
                tooltip = GeoJsonTooltip(
                    fields=colums,
                    aliases=["Delivery country", "Number of shipments"],
                    localize=True
                )
                choropleth.add_child(tooltip)

            

                
                data2=data.groupby(["ZC from"],as_index=False)["PW DSV"].count()
                data2=pd.merge(data2,zip_code,right_on='ZC to',left_on="ZC from")
                data2['count'] = data2.groupby('ZC from')['ZC from'].transform('count')
                data2["PW DSV"]=(data2["PW DSV"])/data2["count"]
                data2=data2.groupby(["nuts0"],as_index=False)["PW DSV"].sum()
                merge2=pd.merge(levl0,data2,right_on="nuts0" ,left_on="ISO2",how="right")


                for k in range (len(merge2)):
                    lat=merge2["LAT"].iloc[k]
                    lon=merge2["LON"].iloc[k]

                    
                    
                    merge2['radius'] = (merge2['PW DSV'] - merge2['PW DSV'].min()) / (merge2['PW DSV'].max() - merge2['PW DSV'].min()) * (20 - 5) + 5
                    merge2 ['radius']= merge2['radius'].fillna(10)
                         
                    folium.CircleMarker(
                        location=[lat, lon],
                        radius=merge2['radius'].iloc[k],
                        color='None',
                        fill=True,
                        fill_color="red",
                        fill_opacity=1,
                        tooltip=f'Collecting country: {merge2["NAME"].iloc[k]} <br> Shipments: {merge2["PW DSV"].iloc[k]}'
                    ).add_to(m)
            
                

                folium_static(m,width=450, height=325)
                st.write("""
                <span style='font-size: small;'>🔴  Collecting countries &nbsp;&nbsp; 🔷  Delivered countries</span>
                """, unsafe_allow_html=True)
                
            col1,col2,col3,col4 = st.columns([1,1,2,2.5])
            with col2:

                df6=data.groupby(['ZC to']).agg({'ZC from': 'count' ,'kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum'  })
                df6=df6.rename(columns={'ZC from' : 'Number of shipments'})
                df6=df6.sort_values(by="Number of shipments",ascending= False )
                df6=df6.head(10)
                df6=df6.sort_values(by="Number of shipments",ascending= True )
                df6 = df6.reset_index()
                fig = px.bar(df6, y='ZC to', x='Number of shipments', 
                color='Number of shipments', 
                color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
                orientation='h',
                hover_data={'kg': True, 'ldm': True, 'PW DSV': True})  
                fig.update_layout(
                    xaxis_title='Shipments',  
                    yaxis_title='' )
                fig.update_coloraxes(showscale=False)
                st.write("<h5><b>Top 10 delivery</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig, use_container_width=True)
                

            with col1:
                df7=data.groupby(['ZC from']).agg({'ZC to': 'count' ,'kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum'  })
                df7=df7.rename(columns={'ZC to' : 'Number of shipments'})
                df7=df7.sort_values(by="Number of shipments",ascending= False )
                df7=df7.head(10)
                df7=df7.sort_values(by="Number of shipments",ascending= True )
                df7 = df7.reset_index()
                fig = px.bar(df7, y='ZC from', x='Number of shipments', 
                color='Number of shipments', 
                color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
                orientation='h',
                hover_data={'kg': True, 'ldm': True})  
                fig.update_layout(
                xaxis_title='Shipments',  
                yaxis_title='' )
                fig.update_coloraxes(showscale=False)
                st.write("<h5><b>Top 10 collection</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig, use_container_width=True)
                
            
            with col3:
                df3=data.groupby(['ZC from','ZC to']).agg({'PW DSV': 'count' ,'kg': 'sum', 'ldm': 'sum'  })
                df3=df3.rename(columns={'PW DSV' : 'Number of shipments'})
                df3=df3.sort_values(by="Number of shipments",ascending=False )
                # df3["%"]=df3["PW DSV"]/(sum(df3["PW DSV"]))*100
                # df3["cum"]=df3["%"].cumsum()
                df3=df3.head(10)
                df3 = df3.reset_index()
                df3["From-to"]=  df3['ZC from'] + ' - ' + df3['ZC to']
                df3.index=df3["From-to"].tolist()
                fig = px.bar(df3, y='Number of shipments', x='From-to',color='Number of shipments', color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'])
                    
                fig.update_layout(
                    xaxis_title='',  
                    yaxis_title='' )
                fig.update_coloraxes(showscale=False)
                st.write("<h5><b>Top 10 main lanes</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig, use_container_width=True)
                
            with col4:
                data['Month'] = pd.to_datetime(data['Date']).dt.to_period('M')
                df4=data.groupby('Month').agg({'Date': 'count' ,'kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum' }).reset_index()
                df4=df4.rename(columns={'Date': 'Shipments'})
                df4['Month'] = df4['Month'].astype(str)
                fig_ship = px.line(df4, x='Month', y='Shipments', markers=True , line_shape='spline')
                fig_ship.update_layout( 
                    xaxis_title='',
                    yaxis_title='Shipments',
                    xaxis={'type': 'category', 'categoryorder': 'array', 'categoryarray': df4['Month'],"showgrid":True})
                st.write("<h5><b>Seasonality</b></h5>", unsafe_allow_html=True)
                st.plotly_chart(fig_ship,use_container_width=True)
                


            col1,col2,col3=st.columns([1.5,1,1.5])

            with col1:
                df4=data.groupby(["Product"]).agg({'ZC to': 'count' ,'kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum'  })
                df4=df4.rename(columns={'ZC to' : 'Shipments'})
                df4 = df4.applymap(lambda x: int(x) if isinstance(x, (int, float)) else x)
                df4=df4.T
                df4["Total"]=df4.sum(axis=1)
                df4 = df4.reset_index()
                df4=df4.rename(columns={'index':"Type"})
                df7=df4.set_index('Type')
                
                df7=df7.applymap(lambda x: '{:,.0f}'.format(x).replace(',', ' ') if pd.notna(x) and isinstance(x, (int, float)) else x)
                st.dataframe(df7)
            with col2:
                df5=data.pivot_table(index="Way", columns="Product", values="PW DSV",aggfunc="count")
                df5["Total"]=df5.sum(axis=1)
                df5= df5.fillna(0)
                
                df5=df5.applymap(lambda x: '{:,.0f}'.format(x).replace(',', ' ') if pd.notna(x) and isinstance(x, (int, float)) else x)
            

            

                st.write(df5)

            


            


elif st.session_state.selected == "Shipment Profile":
        
        data = st.session_state.processed_data
        col1,col2=st.columns([1,7],gap="large")         
        with col1:
            
            if "Category1" in data.columns and "Category2" in data.columns:
                data=filters_with_categories(data)
                
            else:     
                data=apply_filters(data) 
        with col2:
            def can_convert_to_int(value):
                try:
                    int(value)  # Try to convert to integer
                    return True
                except (ValueError, TypeError):
                    return False
            st.header("Shipment Profile")
            st.write(f"***Point of view from {st.session_state.selected_dsv_country} |  Parcel ≤ {st.session_state.selected_parcel} kg | GRP ≤ {st.session_state.selected_grp} kg | LTL ≤ {st.session_state.selected_ltl} kg | FTL > {st.session_state.selected_ltl} kg | Ratios {st.session_state.pw_cbm}/cbm & {st.session_state.pw_ldm}/ldm  {st.session_state.factor_phrase}.***")
            data.dropna(subset=['Bracket'])
            data=data[data['Bracket'].apply(can_convert_to_int)]
            
             
            
            data['Bracket'] = data['Bracket'].astype(str)
            pivot=pd.pivot_table(data,values="PW DSV",index=["Cntry from","ZC from","Cntry to","ZC to"],columns="Bracket",aggfunc="count")
            pivot["total"]=pivot.sum(axis=1)
            pivot["%"]=round(pivot["total"] /( pivot['total'].sum())*100,2)
            pivot.fillna(0, inplace=True)
            pivot=pivot[pivot['total']!=0]

            
            
            # st.dataframe(pivot,width=1000,height=1000)
            st.write('**Total**')
            bracket_columns = [col for col in pivot.columns if col not in ['total', '%']]
            
            pivot[bracket_columns] = pivot[bracket_columns].astype(int)

            
            sorted_bracket_columns = sorted(bracket_columns, key=lambda x: int(x) if x.isdigit() else float('inf'))
            pivot = pivot[sorted_bracket_columns + ['total', '%'] ]
            pivot= pivot.reset_index()
            columns_order = ['%','Cntry from', 'ZC from', 'Cntry to', 'ZC to'] + [col for col in pivot.columns if col not in ['Cntry from', 'ZC from', 'Cntry to', 'ZC to', '%', 'total']]+['total']
            pivot = pivot[columns_order]

            total_row = pivot.sum(numeric_only=True).to_dict()
            total_row.update({'Cntry from': 'Total', 'ZC from': '', 'Cntry to': '', 'ZC to': '', '%': pivot['%'].sum(), 'total': pivot['total'].sum()})
            total_row=pd.DataFrame([total_row])
            
            
            cols = ['%','Cntry from', 'ZC from', 'Cntry to', 'ZC to'] + [col for col in pivot.columns if col not in ['Cntry from', 'ZC from', 'Cntry to', 'ZC to', '%', 'total']]+['total']
            total_row= total_row[cols]
            total_row.drop(columns=['ZC from', 'Cntry to', 'ZC to',"%"], inplace=True)
            percentage=pivot.sum(numeric_only=True).to_dict()
            percentage.update({'Cntry from': '%', 'ZC from': '', 'Cntry to': '', 'ZC to': '', '%': pivot['%'].sum(), 'total': pivot['total'].sum()})
            percentage=pd.DataFrame([percentage])
            
            percentage= percentage[cols]
            percentage.drop(columns=['ZC from', 'Cntry to', 'ZC to',"%"], inplace=True)
            def custom_operation(x):
                if isinstance(x, (int, float)):
                    return round((x / total_row["total"].iloc[0]) * 100, 1)
                return x
            # percentage=(percentage/(total_row["total"].iloc[0])*100).round(1)
            percentage=percentage.applymap(custom_operation)
            
            total_row= pd.concat([percentage,total_row])
            # pivot = pd.concat([pivot, pd.DataFrame([total_row])], ignore_index=True)
           
                
            min_value = pivot['%'].min()
            max_value = pivot['%'].max()
            min_total = pivot['total'].min()
            max_total = pivot['total'].max()
            min_value1 = percentage.drop(columns=["Cntry from", "total"]).min().min()
            max_value1 = percentage.drop(columns=["Cntry from", "total"]).max().max()
            
            
            
            
##################
            gb = GridOptionsBuilder.from_dataframe(pivot)
            

            

            gb.configure_default_column(floatingFilter=True, resizable=False)
            gb.configure_column('Cntry from', headerName="Cntry from", filter="agSetColumnFilter", minWidth=98, maxWidth=300)
            gb.configure_column('ZC from', headerName="ZC from", filter="agSetColumnFilter", minWidth=90, maxWidth=300)
            gb.configure_column('Cntry to', headerName="Cntry to", filter="agSetColumnFilter", minWidth=85, maxWidth=300)
            gb.configure_column('ZC to', headerName="ZC to", filter="agSetColumnFilter", minWidth=80, maxWidth=300)
            gb.configure_column('total', headerName="Total", filter="agNumberColumnFilter", minWidth=70, maxWidth=150)
            gb.configure_column('%', headerName="%", filter="agNumberColumnFilter", minWidth=70, maxWidth=150)

            jscode = JsCode(f"""
                function(params) {{
                    var value = params.value;
                    var minValue = {min_value};
                    var maxValue = {max_value};
                    var normalizedValue = (value - minValue) / (maxValue - minValue) * 100;
                    var color;
                    
                    if (normalizedValue >= 75) {{
                        color = 'green';
                    }} else if (normalizedValue >= 50) {{
                        color = 'orange';
                    }} else {{
                        color = 'red';
                    }}

                    return {{
                        'background': 'linear-gradient(90deg, ' + color + ' ' + normalizedValue + '%,' + ' transparent ' + normalizedValue + '%)',
                        'color': 'black'
                    }};
                }}
            """)
            jscode_total = JsCode("""
                function(params) {
                    var value = params.value;
                    var color = value >= 80 ? '#267326' : // Dark Green
                                value >= 60 ? '#39ac39' : // Medium Green
                                value >= 30 ? '#79d279' : // Medium Green
                                '#d9f2d9'; // Light Green
                    return {
                        'backgroundColor': color,
                        'color': 'black'
                    };
                }
                """)
            jscode1 = JsCode(f"""
                function(params) {{
                    var rowIndex = params.node.rowIndex;
                    var colId = params.column.colId;
                    var value = params.value;
                    var minValue = {min_value1};
                    var maxValue = {max_value1};
                    var normalizedValue = (value - minValue) / (maxValue - minValue) * 100;
                    var color;

                    if (rowIndex === 0 && colId !== 'Total' && colId !== '%') {{
                        if (normalizedValue >= 80) {{
                            color = 'green';
                        }} else if (normalizedValue >= 30) {{
                            color = 'orange';
                        }} else {{
                            color = 'red';
                        }}
                        return {{
                            'background': 'linear-gradient(90deg, ' + color + ' ' + normalizedValue + '%,' + ' transparent ' + normalizedValue + '%)',
                            'color': 'black'
                        }};
                    }}
                    return {{
                        'color': 'black'
                    }};
                }}
            """)
            gb1 = GridOptionsBuilder.from_dataframe(total_row)
            gb1.configure_default_column( resizable=False)
            gb1.configure_column('Cntry from', headerName="Bracket", minWidth=429, maxWidth=429)
            gb.configure_column('%', cellStyle=jscode)
            # gb.configure_column('total', cellStyle=jscode_total)
            for col in total_row.columns:
                if col not in ["Cntry from","total"]:
                    gb1.configure_column(col, cellStyle=jscode1)
            
            gridOptions = gb.build()
            gridOptions1 = gb1.build()
            
            
            AgGrid(total_row,height=91,fit_columns_on_grid_load=True ,allow_unsafe_jscode=True,gridOptions=gridOptions1)
            # Display the AgGrid table in Streamlit
           
            st.write('**Pivot table**')
            st.write("Left filters will apply on the table, you can also: use the filters on top of the table, sort the columns, move the columns, ...")
            
            response = AgGrid(
                pivot,
                gridOptions=gridOptions,
                update_mode=GridUpdateMode.NO_UPDATE,
                enable_enterprise_modules=True,
                allow_unsafe_jscode=True,
                fit_columns_on_grid_load=True,
                height=500
                 
            )
  
elif st.session_state.selected == "Collection Analysis":
    
    data = st.session_state.processed_data
    col1,col2=st.columns([1,7],gap="large")         
    with col1:
        if "Category1" in data.columns and "Category2" in data.columns:
            data=filters_with_categories(data)
            
        else:     
            data=apply_filters(data) 
    with col2:
        st.header('Collection Analysis') 
        st.write(f"***Point of view from {st.session_state.selected_dsv_country} |  Parcel ≤ {st.session_state.selected_parcel} kg | GRP ≤ {st.session_state.selected_grp} kg | LTL ≤ {st.session_state.selected_ltl} kg | FTL > {st.session_state.selected_ltl} kg | Ratios {st.session_state.pw_cbm}/cbm & {st.session_state.pw_ldm}/ldm  {st.session_state.factor_phrase}.***")       
        col1,col2,col3=st.columns([1.3,1.1,1.6],gap='medium')
        with col1:
        
            
            df1=data.groupby('Date').agg({'Date': 'count' ,'kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum' })
            df1=df1.rename(columns={'Date': 'Shipments'})
            df6=df1.applymap(lambda x: '{:,.0f}'.format(x).replace(',', ' ') if pd.notna(x) and isinstance(x, (int, float)) else x) 
            st.dataframe(df6,use_container_width=True,height=510)

            df22=df1.reset_index()
            df22['Date'] = pd.to_datetime(df22['Date'])
            df22['day'] = df22["Date"].dt.day_name()
            df22=df22.groupby('day').agg({'Date': 'count' ,'Shipments':'sum','kg': 'sum', 'ldm': 'sum', 'PW DSV': 'sum' })
            df22=df22.rename(columns={'Date':'Collection'})
            df22=df22.reset_index()
            day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
            df22 = df22.set_index('day').reindex(day_order, fill_value=0).reset_index()
            df22['day'] = pd.Categorical(df22['day'], categories=day_order, ordered=True)
            df22 = df22.sort_values('day')
            
        with col2:
            df2=df1.sum(numeric_only=True).to_frame().T
            df2.index=["Total"]
            df2=df2.applymap(lambda x: '{:,.0f}'.format(x).replace(',', ' ') if pd.notna(x) and isinstance(x, (int, float)) else x)
            st.dataframe(df2,use_container_width=True)
            df7=df1.describe()
            df1.reset_index(inplace=True)
            df1['Date'] = pd.to_datetime(df1['Date'])
            df1['Is_Working_Day'] = df1['Date'].dt.weekday < 5
            num_working_days = df1['Is_Working_Day'].sum()
            total_days = len(df1)
            average_working_days = num_working_days / total_days *5
            df8= pd.DataFrame({
                'Total collection days': [total_days],
                'Working days': [num_working_days],
                'Collection per week': [average_working_days]})
            df8.set_index('Total collection days', inplace=True)
            st.dataframe(df8,use_container_width=True)

            st.dataframe(df7)

            with col3:
                data= data.dropna(subset=['ldm'])
                df3=data.groupby('Date').agg({ 'ldm': 'sum' })
                
                def collection(x):
                    if x <= 0.5:
                        return 0.5
                    elif x <= 1:
                        return 1
                    elif x <= 2:
                        return 2
                    elif x <= 3:
                        return 3
                    elif x <= 4:
                        return 4
                    elif x <= 5:
                        return 5
                    elif x <= 6:
                        return 6
                    elif x <= 7:
                        return 7
                    elif x <= 8:
                        return 8
                    elif x <= 9:
                        return 9
                    elif x <= 10:
                        return 10
                    elif 10<x <= 13.6:
                        return "FTL"
                    elif x > 13.6:
                        return "sup.FTL"
                    
                df3['LDM'] = df3['ldm'].apply(collection)
                
                df3=df3.groupby('LDM').agg({'ldm': 'count' })
                
                df3=df3.rename(columns={'ldm': 'collection'})
                df3=df3.reset_index()
                df4=pd.DataFrame()
                fig=px.bar(df3,x="LDM",y="collection",color="LDM", color_discrete_sequence=['#A9BCE2','#002664','#5D7AB5'], text='collection')
                fig.update_layout(
                                title="Number of collections per ldm  ",
                                xaxis_title='',     
                                yaxis_title='',
                                xaxis=dict(type='category'),
                                showlegend=False,
                            height=280)
                st.plotly_chart(fig)

   

                fig=px.bar(df22,x="day",y="ldm",color="day", color_discrete_sequence=['#002664','#5D7AB5','#A9BCE2'])
                fig.update_layout(
                            title="ldm per week day ",
                            xaxis_title='',     
                            yaxis_title='',
                            showlegend=False,
                            height=300)
                fig.update_coloraxes(showscale=False)
                st.plotly_chart(fig)
            
        data['Month'] = pd.to_datetime(data['Date']).dt.to_period('M')
        df4=data.groupby('Month').agg({'kg': 'count' }).reset_index()
        df4=df4.rename(columns={'kg': 'Collection'})
        df4['Month'] = df4['Month'].astype(str)
        fig_ship = px.line(df4, x='Month', y='Collection', markers=True, line_shape='spline')
        fig_ship.update_layout(
                    title='Seasonality of collection per month', 
                    xaxis_title='Month',
                    yaxis_title='Number of collection',
                    height=290
                )
        st.plotly_chart(fig_ship,use_container_width=False)

           
#########################
           

            
          

                      ##################################################################################################
# Page for Map
elif st.session_state.selected == "Maps":
    
    st.write("**Filters option**")
    
    data = st.session_state.processed_data

    col1,col2=st.columns([1,7],gap="large")         
    with col1:
        if not data.empty:
                
                selected_branch = st.multiselect('Select branch', options=data['Branch'].unique())
                selected_level = st.selectbox('Select a level', ["country level", "Nuts1", "Nuts2", "Nuts3"])
                produit= st.multiselect('Select type of product', options=data['Product'].unique())
                selected_variable=st.selectbox('Select the variable to see the map based on', options=["Number of shipments","PW DSV","kg","ldm","m3"])
                if data['Date'].isna().sum() == 0:
                    data['Date'] = pd.to_datetime(data['Date'])

                    first_date = data['Date'].min()
                    last_date = data['Date'].max()

                    start_date = st.date_input(
                        "Select start date",
                        value=first_date,
                        min_value=first_date,
                        max_value=last_date,
                        format="DD.MM.YYYY"
                        
                    )

                    end_date = st.date_input(
                        "Select end date",
                        value=last_date,
                        min_value=start_date,
                        max_value=last_date,
                        format="DD.MM.YYYY"
                    )
                
                
              

                if 'start_date' in locals() and 'end_date' in locals():
                    data=data[(data['Product'].isin(produit) if produit else data['Product'].notnull())&
                    (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())&
                    (data['Date'] >= pd.to_datetime(start_date)) & (data['Date'] <= pd.to_datetime(end_date))]
                else:
                    data=data[(data['Product'].isin(produit) if produit else data['Product'].notnull())&
                    (data['Branch'].isin(selected_branch) if selected_branch else data['Branch'].notnull())]   
                m= folium.Map(location=[54.5260,15.2551],zoom_start=4,width='100%', control_scale=True)
                m1= folium.Map(location=[54.5260,15.2551],zoom_start=4,width='100%', control_scale=True)
                if st.checkbox('show DSV branches'):
                    categories=dsv["Country"].unique().tolist()
                    for k in range(len(dsv)):
                        location = dsv["lat"].iloc[k],dsv["lon"].iloc[k]
                        html= f"<span style='font-family:Arial;''><span style='color:#002664;'><b>{ dsv['Office Name'].iloc[k]} <br> { dsv['ZC'].iloc[k]}  </b> </span><br> <br>  &#9743; {dsv['Phone'].iloc[k]} <br> <span style='font-size:12px'> &#128343; </span> {dsv['Opening hours'].iloc[k]} </span>"
                        iframe = branca.element.IFrame(html=html, width=250, height=140)

                        folium.Marker(
                        location=location,
                        tags=[dsv["Country"].iloc[k]],
                        icon=folium.Icon(color='darkblue',icon_color='White', icon="info-sign"),
                        tooltip=dsv["Office Name"].iloc[k] + " , click for more information",
                        popup = folium.Popup(iframe, max_width=500)).add_to(m)
                     

                        folium.Marker(
                        location=location,
                        tags=[dsv["Country"].iloc[k]],
                        icon=folium.Icon(color='darkblue',icon_color='White', icon="info-sign"),
                        tooltip=dsv["Office Name"].iloc[k] + " , click for more information",
                        popup = folium.Popup(iframe, max_width=500)).add_to(m1)
                    TagFilterButton(categories).add_to(m)
                    TagFilterButton(categories).add_to(m1) 

        if selected_variable=="Number of shipments":

            df7=data.groupby(['ZC from']).agg({'PW DSV': 'count' ,'kg': 'sum', 'ldm': 'sum','m3':'sum'  })
            df7=df7.rename(columns={'PW DSV' : 'Number of shipments'})
            df7=df7.sort_values(by="Number of shipments",ascending= False )
            df7=df7.head(10)
            df7=df7.sort_values(by="Number of shipments",ascending= True )
            df7 = df7.reset_index()
            fig = px.bar(df7, y='ZC from', x='Number of shipments', 
            color='Number of shipments', 
            color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
            orientation='h',
            hover_data={'kg': True, 'ldm': True})  
            fig.update_layout(
            xaxis_title='Shipments',  
            yaxis_title='',
            height=300,
            title="Top 10 collections ZC" )
            fig.update_coloraxes(showscale=False)
            st.plotly_chart(fig, height=50)

            df6=data.groupby(['ZC to']).agg({'PW DSV': 'count' ,'kg': 'sum', 'ldm': 'sum','m3':'sum'  })
            df6=df6.rename(columns={'PW DSV' : 'Number of shipments'})
            df6=df6.sort_values(by="Number of shipments",ascending= False )
            
            
            df6=df6.head(10)
            df6=df6.sort_values(by="Number of shipments",ascending= True )
            df6 = df6.reset_index()
            fig = px.bar(df6, y='ZC to', x='Number of shipments', 
            color='Number of shipments', 
            color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
            orientation='h',
            hover_data={'kg': True, 'ldm': True})  
            fig.update_layout(
                xaxis_title='Shipments',  
                yaxis_title='',
                height=300,
                title= "Top 10 deliveries ZC")
            fig.update_coloraxes(showscale=False)
            
            st.plotly_chart(fig, use_container_width=True)
        else:

            df7=data.dropna(subset=[selected_variable],how='any')
            df7=df7.groupby(['ZC from']).agg({'PW DSV': 'sum' ,'kg': 'sum', 'ldm': 'sum','m3':'sum'  })
            df7=df7.sort_values(by=selected_variable,ascending= False )
            df7=df7.head(10)
            df7=df7.sort_values(by=selected_variable,ascending= True )
            df7 = df7.reset_index()
            fig = px.bar(df7, y='ZC from', x=selected_variable, 
            color=selected_variable, 
            color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
            orientation='h',
            hover_data={'kg': True, 'ldm': True})  
            fig.update_layout(
            xaxis_title=f'Total {selected_variable}',  
            yaxis_title='',
            height=300,
            title="Top 10 collections ZC" )
            fig.update_coloraxes(showscale=False)
            st.plotly_chart(fig, height=50)


            df6=data.dropna(subset=[selected_variable],how='any')
            df6=df6.groupby(['ZC to']).agg({'PW DSV': 'sum' ,'kg': 'sum', 'ldm': 'sum','m3':'sum'  })
            df6=df6.sort_values(by=selected_variable,ascending= False )
            df6=df6.head(10)
            df6=df6.sort_values(by=selected_variable,ascending= True )
            df6 = df6.reset_index()
            fig = px.bar(df6, y='ZC to', x=selected_variable, 
            color=selected_variable, 
            color_continuous_scale=['#A9BCE2','#5D7AB5','#002664'],
            orientation='h',
            hover_data={'kg': True, 'ldm': True})  
            fig.update_layout(
                xaxis_title=f'Total {selected_variable}',  
                yaxis_title='',
                height=300,
                title= "Top 10 deliveries ZC")
            fig.update_coloraxes(showscale=False)
            
            st.plotly_chart(fig, use_container_width=True)
        
        
                    

        
   
    with col2:
        st.header('Map')
        st.write(f"***Point of view from {st.session_state.selected_dsv_country} |  Parcel ≤ {st.session_state.selected_parcel} kg | GRP ≤ {st.session_state.selected_grp} kg | LTL ≤ {st.session_state.selected_ltl} kg | FTL > {st.session_state.selected_ltl} kg | Ratios {st.session_state.pw_cbm}/cbm & {st.session_state.pw_ldm}/ldm  {st.session_state.factor_phrase}.***")
        tab1,tab2=st.tabs(["Outbound from a single ZC","Inbound to a single ZC"])
        with tab1:
            col3,col4=st.columns([1,9])
            with col3:
                ship_from = sorted(data['ZC from'].dropna().unique().tolist())
                selected_country = st.selectbox('Select Shipments from:', ship_from)
            with col4:
                st.empty()

            data_to = data[(data['ZC from'] == selected_country)]
            
            if selected_variable!="Number of shipments":
                data_to=data_to.groupby(["ZC to"],as_index=False)[selected_variable].sum()
                selected_variable_column=selected_variable
            else:
                selected_variable="PW DSV"
                data_to=data_to.groupby(["ZC to"],as_index=False)[selected_variable].count()
                selected_variable_column="Number of shipments"
                
            data_to=pd.merge(data_to,zip_code,on='ZC to',how="left")
            data_to['count'] = data_to.groupby('ZC to')['ZC to'].transform('count')
            data_to[selected_variable]=(data_to[selected_variable])/data_to["count"]

    
        

            if selected_level== "country level":
                data_to=data_to.groupby(["nuts0"],as_index=False)[selected_variable].sum()
                merge=pd.merge(levl0,data_to,right_on="nuts0" ,left_on="ISO2",how="right")
                colums=["nuts0",selected_variable]
                key="properties.ISO2"
                field=["NAME",selected_variable]
                alias=["To : ", f"{selected_variable_column}: "]
                
            
            elif selected_level== "Nuts1":
                data_to=data_to.groupby(["nuts1"],as_index=False)[selected_variable].sum()
                merge=pd.merge(levl1,data_to,right_on="nuts1" ,left_on="NUTS_ID",how="right")
                colums=["nuts1",selected_variable]
                key="properties.NUTS_ID"
                field=["NUTS_NAME","NUTS_ID",selected_variable]
                alias=["To: " ,"NUTS_ID: ",  f"{selected_variable_column}: "]
            elif selected_level== "Nuts2":
                data_to=data_to.groupby(["nuts2"],as_index=False)[selected_variable].sum()
                merge=pd.merge(levl2,data_to,right_on="nuts2" ,left_on="NUTS_ID",how="right")
                colums=["nuts2",selected_variable]
                key="properties.NUTS_ID"
                field=["NUTS_NAME","NUTS_ID",selected_variable]
                alias=["To: " ,"NUTS_ID: ",  f"{selected_variable_column}: "]
            elif selected_level== "Nuts3":
                data_to=data_to.groupby(["NUTS3"],as_index=False)[selected_variable].sum()
                merge=pd.merge(levl3,data_to,right_on="NUTS3" ,left_on="NUTS_ID",how="right")
                colums=["NUTS3",selected_variable]
                key="properties.NUTS_ID"
                field=["NUTS_NAME","NUTS_ID",selected_variable]
                alias=["To: " ,"NUTS_ID: ",  f"{selected_variable_column}: "]

    
            
            choropleth=folium.Choropleth(
            geo_data=merge,
            name="choropleth",
            data=merge,
            columns=colums,
            key_on=key,
            fill_color='OrRd',
            fill_opacity=0.7,
            legend=True,
            highlight=True,                   
            ).geojson.add_to(m)

            folium.features.GeoJson(
                            data=merge,
                            name='test',
                            smooth_factor=2,
                            style_function=lambda x: {'color':'black','fillColor':'transparent','weight':0.5},
                            tooltip=folium.features.GeoJsonTooltip(
                                fields=field,    
                                aliases=alias,  
                                localize=True,
                                sticky=False,
                                labels=True,
                                style="""
                                    background-color: #F0EFEF;
                                    border: 2px solid black;
                                    border-radius: 3px;
                                    box-shadow: 3px;
                                """,
                                max_width=800,),
                                    highlight_function=lambda x: {'weight':3,'fillColor':'grey'},
                                ).add_to(choropleth)
        
            if selected_country == selected_country :
                try:
                    nuts_of_the_ZC=zip_code["NUTS3"].loc[zip_code["ZC to"]==selected_country]
                    polygon = levl3['geometry'].loc[levl3["NUTS_ID"]==(nuts_of_the_ZC.tolist()[0])]
                    centroid = polygon.centroid
                    long=centroid.x.tolist()[0]
                    lat=centroid.y.tolist()[0]
                    folium.Marker([lat, long],icon=folium.Icon(color='red', ), tooltip=f" From {selected_country}").add_to(m)
                except Exception :
                    print ("error")
            Fullscreen(position="topleft").add_to(m)
            html_string = m.get_root().render()
            
            folium_static(m, height=700,width=1200)
            
            def create_download_button(html_string, filename):
                    b64 = base64.b64encode(html_string.encode()).decode()
                    href = f'<a href="data:text/html;base64,{b64}" download="{filename}">Download this map </a>'
                    return href
            st.markdown(create_download_button(html_string, "map.html"), unsafe_allow_html=True)

            with tab2:
                col3,col4=st.columns([1,9])
                with col3:
                    ship_to=sorted(data['ZC to'].dropna().unique().tolist())
                    selected_country_to = st.selectbox('Select Shipments to:', ship_to)
                with col4:
                    st.empty()
                    
                data_to = data[(data['ZC to'] == selected_country_to)]

                if selected_variable!="Number of shipments":

                    data_to=data_to.groupby(["ZC from"],as_index=False)[selected_variable].sum()
                    selected_variable_column=selected_variable
                else:
                    selected_variable="PW DSV"
                    data_to=data_to.groupby(["ZC from"],as_index=False)[selected_variable].count()
                    selected_variable_column="Number of shipments"

                data_to=pd.merge(data_to,zip_code,right_on='ZC to',left_on="ZC from",how="left")
                data_to['count'] = data_to.groupby('ZC from')['ZC from'].transform('count')
                data_to[selected_variable]=(data_to[selected_variable])/data_to["count"]

        
            

                if selected_level== "country level":
                    data_to=data_to.groupby(["nuts0"],as_index=False)[selected_variable].sum()
                    merge_to=pd.merge(levl0,data_to,right_on="nuts0" ,left_on="ISO2",how="right")
                    colums=["nuts0",selected_variable]
                    key="properties.ISO2"
                    field=["NAME",selected_variable]
                    alias=["From : ", f"{selected_variable_column} : "]
                    
                
                elif selected_level== "Nuts1":
                    data_to=data_to.groupby(["nuts1"],as_index=False)[selected_variable].sum()
                    merge_to=pd.merge(levl1,data_to,right_on="nuts1" ,left_on="NUTS_ID",how="right")
                    colums=["nuts1",selected_variable]
                    key="properties.NUTS_ID"
                    field=["NUTS_NAME","NUTS_ID",selected_variable]
                    alias=["From: " ,"NUTS_ID: ",  f"{selected_variable_column} : "]
                elif selected_level== "Nuts2":
                    data_to=data_to.groupby(["nuts2"],as_index=False)[selected_variable].sum()
                    merge_to=pd.merge(levl2,data_to,right_on="nuts2" ,left_on="NUTS_ID",how="right")
                    colums=["nuts2",selected_variable]
                    key="properties.NUTS_ID"
                    field=["NUTS_NAME","NUTS_ID",selected_variable]
                    alias=["From: " ,"NUTS_ID: ",  f"{selected_variable_column} : "]
                elif selected_level== "Nuts3":
                    data_to=data_to.groupby(["NUTS3"],as_index=False)[selected_variable].sum()
                    merge_to=pd.merge(levl3,data_to,right_on="NUTS3" ,left_on="NUTS_ID",how="right")
                    colums=["NUTS3",selected_variable]
                    key="properties.NUTS_ID"
                    field=["NUTS_NAME","NUTS_ID",selected_variable]
                    alias=["From: " ,"NUTS_ID: ",  f"{selected_variable_column} : "]
        
                
                choropleth1=folium.Choropleth(
                geo_data=merge_to,
                name="choropleth",
                data=merge_to,
                columns=colums,
                key_on=key,
                fill_color='Blues',
                fill_opacity=0.7,
                legend=True,
                highlight=True,                   
                ).geojson.add_to(m1)

                folium.features.GeoJson(
                                data=merge_to,
                                name='test',
                                smooth_factor=2,
                                style_function=lambda x: {'color':'black','fillColor':'transparent','weight':0.5},
                                tooltip=folium.features.GeoJsonTooltip(
                                    fields=field,    
                                    aliases=alias,  
                                    localize=True,
                                    sticky=False,
                                    labels=True,
                                    style="""
                                        background-color: #F0EFEF;
                                        border: 2px solid black;
                                        border-radius: 3px;
                                        box-shadow: 3px;
                                    """,
                                    max_width=800,),
                                        highlight_function=lambda x: {'weight':3,'fillColor':'grey'},
                                    ).add_to(choropleth1)
                if selected_country_to == selected_country_to :
                    try:
                        nuts_of_the_ZC=zip_code["NUTS3"].loc[zip_code["ZC to"]==selected_country_to]
                        polygon = levl3['geometry'].loc[levl3["NUTS_ID"]==(nuts_of_the_ZC.tolist()[0])]
                        centroid = polygon.centroid
                        long=centroid.x.tolist()[0]
                        lat=centroid.y.tolist()[0]
                        folium.Marker([lat, long],icon=folium.Icon(color='blue', ), tooltip=f" To {selected_country_to}").add_to(m1)
                    except Exception :
                        print ("error")
                Fullscreen(position="topleft").add_to(m1)
                html_string = m1.get_root().render()
                
                folium_static(m1, height=700,width=1200)
                
                def create_download_button(html_string, filename):
                        b64 = base64.b64encode(html_string.encode()).decode()
                        href = f'<a href="data:text/html;base64,{b64}" download="{filename}">Download this the map </a>'
                        return href
                st.markdown(create_download_button(html_string, "map.html"), unsafe_allow_html=True)
            
elif st.session_state.selected == "Regularity Detector":
    data = st.session_state.processed_data
    col1,col2=st.columns([1,7],gap="large")         
    with col1:
                if "Category1" in data.columns and "Category2" in data.columns:
                    data=filters_with_categories(data)
                    
                else:     
                  data=apply_filters(data) 
    with col2:
        st.header('Regularity Detector') 
        st.write(f"***Point of view from {st.session_state.selected_dsv_country} |  Parcel ≤ {st.session_state.selected_parcel} kg | GRP ≤ {st.session_state.selected_grp} kg | LTL ≤ {st.session_state.selected_ltl} kg | FTL > {st.session_state.selected_ltl} kg | Ratios {st.session_state.pw_cbm}/cbm & {st.session_state.pw_ldm}/ldm  {st.session_state.factor_phrase}.***")
        data = data[data['Date'].notna()]
        data['Date'] = pd.to_datetime(data['Date'])
        data['day'] = data["Date"].dt.day_name()
           
        
        data['lane'] = data['ZC from'].astype(str) + ' - ' + data['ZC to'].astype(str)       
       
        days_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
        data['day'] = pd.Categorical(data['day'], categories=days_order, ordered=True)
        
        ldm_mean = data.groupby(['lane', 'Branch', 'Product', 'Way'])['ldm'].mean().reset_index()
        df=pd.pivot_table(data,values="PW DSV",index=["lane",'Branch','Product','Way'],columns="day",aggfunc="count")
        

        df['Total'] = df.sum(axis=1) 
        
        df = df.reset_index()
        
        ldm_mean.columns = ['lane', 'Branch', 'Product', 'Way', 'Average LDM']
        

        df=df.dropna(how='all', subset=df.columns[4:])
        df=df[df['Total'] != 0]
        df["Average ldm"]=ldm_mean['Average LDM']
        new_column_order = ['lane', 'Branch', 'Product', 'Way', 'Average ldm'] + [col for col in df.columns if col not in ['lane', 'Branch', 'Product', 'Way', 'Average ldm']]
        df=df[new_column_order]
        
        # st.write(df)
        if not df.empty:
            gb = GridOptionsBuilder.from_dataframe(df)
            cell_style_jscode = JsCode(  '''
            function(params) {
                const totalValue = params.data.Total;  // Access the row's total value
                const cellValue = params.value;  // Access the current cell's value

                // Apply percentage-based styling only if the column is not 'Total'
                if (params.colDef.field !== 'Total' && totalValue && totalValue !== 0) {  // Exclude 'Total' column
                    const percentage = (cellValue / totalValue) * 100;  // Calculate the percentage

                    if (percentage == 100) {
                        return {'backgroundColor': 'rgb(105, 44, 17)', 'color': 'white'};  
                    } else if (percentage > 92) {
                        return {'backgroundColor': 'rgb(128, 67, 9)', 'color': 'white'};  
                    } else if (percentage > 83) {
                        return {'backgroundColor': 'rgb(150, 89, 0)', 'color': 'white'};  
                    } else if (percentage > 75) {
                        return {'backgroundColor': 'rgb(172, 100, 24)', 'color': 'black'};  
                    } else if (percentage > 67) {
                        return {'backgroundColor': 'rgb(194, 111, 48)', 'color': 'black'}; 
                    } else if (percentage > 58) {
                        return {'backgroundColor': 'rgb(219, 116, 40)', 'color': 'black'};  
                    } else if (percentage > 50) {
                        return {'backgroundColor': 'rgb(244, 121, 32)', 'color': 'black'};  
                    } else if (percentage > 42) {
                        return {'backgroundColor': 'rgb(242, 146, 67)', 'color': 'black'};  
                    } else if (percentage > 33) {
                        return {'backgroundColor': 'rgb(240, 171, 102)', 'color': 'black'};  
                    } else if (percentage > 25) {
                        return {'backgroundColor': 'rgb(243, 181, 125)', 'color': 'black'};  
                    } else if (percentage > 17) {
                        return {'backgroundColor': 'rgb(245, 191, 147)', 'color': 'black'};  
                    } else if (percentage > 8) {
                        return {'backgroundColor': 'rgb(255, 219, 191)', 'color': 'black'};  
                    } else if (percentage <8) {
                        return {'backgroundColor': 'rgb(255, 255, 255)', 'color': 'black'};  
                    }
                } else if (params.colDef.field === 'Total') {
                    return {'backgroundColor': '#fdfefe', 'color': 'black'};  // Default styling for 'Total' column
                }
            }
            ''' )
            gb.configure_column(
                field="Weekly Trend",
                headerName="Weekly Trend",
                cellRenderer='agSparklineCellRenderer',  # Important to use the correct renderer
                cellRendererParams={
                    'sparklineOptions': {
                        'line': {
                            'strokeWidth': 2,  # Optional: Set line width
                            'stroke': 'blue',  # Optional: Set color of the line
                        },
                        'highlightStyle': {
                            'fill': 'red',  # Optional: Color to highlight a data point
                        }
                    }
                }
            )
            # Add 'Weekly Trend' column to DataFrame (required by AG Grid but can remain empty)
            df['Weekly Trend'] = df.apply(lambda row: [
                row['Monday'], row['Tuesday'], row['Wednesday'], row['Thursday'],
                row['Friday'], row['Saturday'], row['Sunday']
            ], axis=1)

            gb.configure_default_column( filter=True, sortable=True, floatingFilter=True) 
            gb.configure_column('Monday', cellStyle=cell_style_jscode, headerName="Mon")
            gb.configure_column('Tuesday', cellStyle=cell_style_jscode, headerName="Tues")
            gb.configure_column('Wednesday', cellStyle=cell_style_jscode, headerName="Wed")
            gb.configure_column('Thursday', cellStyle=cell_style_jscode, headerName="Thurs")
            gb.configure_column('Friday', cellStyle=cell_style_jscode, headerName="Fri")
            gb.configure_column('Saturday', cellStyle=cell_style_jscode, headerName="Sat")
            gb.configure_column('Sunday', cellStyle=cell_style_jscode, headerName="Sun")
            gb.configure_column('Total', cellStyle=cell_style_jscode)

            gb.configure_column('lane', headerName="Lane", filter="agSetColumnFilter", minWidth=124)
            gb.configure_column('Branch', headerName="Branch", filter="agSetColumnFilter", minWidth=120)
            gb.configure_column('Product', headerName="Product", filter="agSetColumnFilter", minWidth=109)
            gb.configure_column('Way', headerName="Way", filter="agSetColumnFilter", minWidth=10) 
            gb.configure_column('Average ldm', headerName="Avg ldm", filter="agSetColumnFilter", minWidth=112)
            gb.configure_column('Weekly Trend', headerName="Weekly Trend",  minWidth=140)  
            grid_options = gb.build()
            
            response = AgGrid(
                    df,
                    gridOptions=grid_options,
                    update_mode=GridUpdateMode.NO_UPDATE,
                    enable_enterprise_modules=True,
                    allow_unsafe_jscode=True,
                    fit_columns_on_grid_load=True,
                    theme="material",
                    height=600      
                )
       
        

            
      


















elif st.session_state.selected == "Document":
    PASSWORD = "Dsv2025+"

    # Session state to track if the user is authenticated
    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False

    def authenticate(password):
        """Authenticate the user by checking the password."""
        if password == PASSWORD:
            st.session_state["authenticated"] = True
            st.success("Authentication successful!")
            st.rerun()
        else:
            st.error("Invalid password. Please try again.")

    # If the user is not authenticated, show the password input
    if not st.session_state["authenticated"]:
        
        col1, col2, col3 = st.columns(3)
        with col2:
            st.title("Authentication Required")
        
        password_input = st.text_input("Enter password", type="password")
        if st.button("Login"):
            authenticate(password_input)
    else:
        
            # def split_text_into_chunks(text, max_tokens=1500):
            #     encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
            #     tokens = encoding.encode(text)
            #     chunks = []
            #     current_chunk = []
            #     current_tokens = 0

            #     for token in tokens:
            #         current_chunk.append(token)
            #         current_tokens += 1
            #         if current_tokens >= max_tokens:
            #             chunks.append(encoding.decode(current_chunk))
            #             current_chunk = []
            #             current_tokens = 0

            #     if current_chunk:
            #         chunks.append(encoding.decode(current_chunk))

                # return chunks

            
            st.header("Upload a file and ask me anything about it!") 
            
            
            with st.sidebar:
                st.title("")
                st.title("")
                openai_api_key = st.text_input("Put your api key in here and press enter")
                
                uploaded_file = st.file_uploader("Upload your file", type=["pdf", "pptx", "docx"])
                if uploaded_file:
                    st.success("PDF uploaded successfully!")

                selected_language=st.selectbox("Change language",options=["English", "French", "Spanish", "Italian"])
                
                if st.button("Logout"):
                    st.session_state["authenticated"] = False
                    st.rerun()
          
            

            # Function to extract text from PDF
            def extract_text_from_pdf(pdf_file):
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                text = ""
                for page_num in range(len(pdf_reader.pages)):
                    text += pdf_reader.pages[page_num].extract_text()
                return text
            def extract_text_from_pptx(pptx_file):
            
                presentation = Presentation(pptx_file)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text_frame.text + "\n"
                return text
            
            def extract_text_from_docx(docx_file):
                document = Document(docx_file)
                text = ""
                for paragraph in document.paragraphs:
                    text += paragraph.text + "\n"
                return text
            
            def extract_text_from_file(file, file_type):
                if file_type == "pdf":
                    return extract_text_from_pdf(file)
                elif file_type == "pptx":
                    return extract_text_from_pptx(file)
                elif file_type == "docx":
                    return extract_text_from_docx(file)
                else:
                    print("Unsupported file type.")
                    return ""
                
            tab1,tab3,tab2= st.tabs(["Summary","Strategy/PPt presentation","Chat"])
            
            document_text = ""
            contxt="Context=You work in an European country for the Road division of a major international supplier of transport and logistics solutions offering a full range of services. The Road division -  which is the top leading road transport provider in Europe with 200 terminals and 300 offices – offers to its customers groupage (GRP) through his fully integrated network, partloads (LTL) and Full Trailer Loads (FTL). Your company is specialized in palletized industrial goods."
            message_placeholder = st.empty()
            if uploaded_file is not None:
                
                file_type = uploaded_file.name.split(".")[-1].lower()
                document_text = extract_text_from_file(uploaded_file, file_type)
                # document_chunks = split_text_into_chunks(document_text)
                
                
                if not openai_api_key:
                            st.info("Please add your OpenAI API key to continue.")
                            st.stop()
            with tab1:
                if uploaded_file:
                    
                        
                    messages = [{'role': 'system', 'content': f'Answer in {selected_language} .in bullet points I want specific answers of this, only if they exist if not just type “no information” :Customer name, Project number/name, Customer sector, Expected number of rounds, deadline to answer the tender, Customer decision date, General customer info, Tender scope, Award strategy, Contract validity, Rate validity, Start date,Payment terms, Palletized goods, Parcel yes if it exist no if not, groupage yes if it exist no if not, LTL yes if it exist no if not, FTL “yes” if it exist “no” if not, Intermodal yes if it exist no if not, Box trailers yes if it exist no if not, Curtain trailers yes if it exist no if not, Mega trailers yes if it exist no if not, Open trailers yes if it exist no if not, Reefer trailers yes if it exist no if not, Jumbo trailers yes if it exist no if not, Temperature controlled yes if it exist no if not, KFF - Keep From Freezing yes if it exist no if not, Tail-lift yes if it exist no if not, ADR yes if it exist no if, Stand trailer yes if it exist no if not, Control Tower yes if it exist no if not and (if yes add “:” and write a description of the control tower needs), Total kg (if mentioned type “:” and write “estimation in euro” then the result of the multiplication of total kg by 0.12), Penalties yes if Financial Penalties are explicitly mentioned no if not (if yes type “:”  and write what are they in the same line), the Fuel Clause Mechanism, the Fuel share in rates, Threshold fuel price, Baseline reference date, Baseline price, Fuel based on, Current Fuel price, Current Fuel surcharge % ,Calculation Date FSC, Go/NoGo:Go if none of the showstopper exist and NoGo if one or more showstopper exist and if no go put “:” and write what are the showstopper if Go put “:” and write that there are no showstoppers found.(showstopper=100% of transports concern thermo trailers (temperature control),100% of transports concern parcels or non-palletized goods or bulk shipments, they concern non-industrial goods (ie living animals),there are financial penalties, palette exchange.) .This is the document content :\n\n{document_text}'}]


                    try:
                        client = OpenAI(api_key=openai_api_key) 
                        response = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=messages)
                        msg = response.choices[0].message.content
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
                        

                    messages_summary = [{"role": "system", "content": f"Answer in {selected_language} I want a summary of this document in one paragraph and then in bullet points i want recommendation of a strategy to win this tender:\n\n{document_text}"}]
                    try:
                        client = OpenAI(api_key=openai_api_key) 
                        response = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=messages_summary)
                        summary_strategy = response.choices[0].message.content
                        split_response = summary_strategy.split("\n\n", 1)
                        summary = split_response[0].strip()  
                        recommendations = split_response[1].strip() if len(split_response) > 1 else ""
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
                        
                        
                    st.subheader("Summary")
                    if summary:    
                        st.write(summary)
                        summary_dataframe = pd.DataFrame({"Summary": [summary]})
                    else:
                        st.warning("Unable to create THE summary, please rerun. If error persist check your document")
                        if st.button("Rerun"):
                            st.rerun()                       

                    

                    messages = [{"role": "system", "content": f"Answer in {selected_language} As an expert in sales and marketing create a PowerPoint to present to the client to win this tender. Each slide should have a clear title and 3 to 5 bullet points for content and cover the needs and requirements available on the document. Keep the language concise and professional, and structure it to engage an audience of transport buyers and logistic managers.\n\n{document_text}"}]
                    try:
                        client = OpenAI(api_key=openai_api_key) 
                        response = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=messages)
                        ppt = response.choices[0].message.content
                        ppt_dataframe = pd.DataFrame({"PPt Presentation": [ppt]})
                    except Exception as e:
                            st.error(f"An error occurred: {e}")
                    


                    col1,col2=st.columns([2,1.5])
                    with col1:

                        st.subheader("Questions")
                        if msg: 
                            lines = [line[1:] for line in msg.split("\n")[:33]]
                            structured_data = [line.split(": ", 2) for line in lines if ": " in line]
                            structured_data=[columns if len(columns)==3 else columns + [""]for columns in structured_data]

                            if structured_data:
                                df = pd.DataFrame(structured_data, columns=["Question", "Answer","Description"])
                                
                                df1=df.set_index("Question")
                                st.table(df1)
                                try:
                                    excel_name=f'profilor RRF-{df["Answer"][0]}-deadline-{df["Answer"][4]}.xlsx'
                                except:
                                    excel_name="Profilor RRF.xlsx"
                            else:
                                st.warning("Unable to create table, please rerun. If error persist check your document")
                                if st.button("Rerun"):
                                    st.rerun()
                        else:
                            st.warning("Unable to create table, please rerun. If error persist check your document")
                            if st.button("Rerun"):
                                st.rerun()



                    with col2:
                        st.subheader("FSC")
                        if msg:  
                            lines_fuel = [line[1:] for line in msg.split("\n")[33:-1]]
                            structured_data_fuel = [line.split(": ", 1) for line in lines_fuel if ": " in line]



                            if structured_data_fuel:
                                df_fuel = pd.DataFrame(structured_data_fuel, columns=["Question", "Answer",])
                                df_fuel1=df_fuel.set_index("Question")
                                st.table(df_fuel1)


                            st.subheader("GO or NoGO") 
                            lines_go = [line[1:] for line in msg.split("\n")[-1:]]
                            structured_data_go = [line.split(": ", 2) for line in lines_go if ": " in line]
                            structured_data_fuel=[columns if len(columns)==3 else columns + [""]for columns in structured_data]
                            if structured_data_go:
                                df_go = pd.DataFrame(structured_data_go, columns=["Question", "Answer","Description"])
                                df_go1=df_go.set_index("Question")
                                st.table(df_go1)
                        else:
                            st.warning("Unable to create table, please rerun. If error persist check your document")
                            if st.button("Rerun"):
                                st.rerun()
                        



                    if not df.empty and not df_fuel.empty and not df_go.empty:
                        def convert_df_to_excel(datas):
                            output = BytesIO()
                            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                                for sheet_name,dataframe in datas.items():
                                    dataframe.to_excel(writer, index=False, sheet_name=sheet_name)

                                    workbook = writer.book
                                    worksheet = writer.sheets[sheet_name]
                                    
                                    # Add custom formats
                                    header_format = workbook.add_format({
                                        'bold': True,
                                        'font_color': 'white',
                                        'bg_color': '#4F81BD',
                                        'align': 'center',
                                        'valign': 'vcenter',
                                        'border': 1
                                    })
                                    cell_format = workbook.add_format({
                                        'text_wrap': True,
                                        'valign': 'top',
                                        'border': 1
                                    })
                                    
                                    # Apply header format
                                    for col_num, value in enumerate(dataframe.columns.values):
                                        worksheet.write(0, col_num, value, header_format)
                                    
                                    # Apply cell format to the data
                                    for row_num in range(1, len(dataframe) + 1):
                                        for col_num in range(len(dataframe.columns)):
                                            worksheet.write(row_num, col_num, dataframe.iloc[row_num - 1, col_num], cell_format)
                                    
                                    # Adjust column widths
                                    max_width = 100
                                    for i, col in enumerate(dataframe.columns):
                                        max_len = max(
                                            dataframe[col].astype(str).map(len).max(),  # Maximum length in column
                                            len(str(col))  # Length of the column header
                                        ) + 2  # Add some padding

                                        adjusted_width = min(max_len, max_width)
                                        worksheet.set_column(i, i, adjusted_width)

                            return output.getvalue()

                    datas={"Summary":summary_dataframe,
                            "Standard questions":df ,
                            "Fuel surchage":df_fuel,
                            "Go no go":df_go,
                            "PPt Presentation": ppt_dataframe}
                    # Add download button for Excel
                    with st.sidebar:
                        excel_data = convert_df_to_excel(datas)
                        st.download_button(
                            label="📥 Download as an Excel File",
                            data=excel_data,
                            file_name=excel_name,
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                else:
                    st.info("Upload a document to see the summary")

            with tab2:
                if 'generated' not in st.session_state:
                    st.session_state['generated'] = []
                if 'past' not in st.session_state:
                    st.session_state['past'] = []
                
                chat_context=st.radio("",options=["Questions about the document","General question"])
                if chat_context=="Questions about the document":
                    context=f"This is the document content:\n\n{document_text}"
                if chat_context=="General question":
                    context="."
                chat_role=st.selectbox("Ask a specialist", options=["Business analyst", "Legal consultant","Sales specialist","CEO", "Not a specialist"])
                with st.expander("See and edit prompt (Press enter to apply)"):
                    editable_prompt = st.text_input("",value=f"Answer as if you are a {chat_role}  {contxt} ")
                



                if "messages" not in st.session_state:
                    st.session_state["messages"] = [{"role": "assistant", "content": "Hello, how can I help you?"}]


                def generate_response(prompt):
                    
                    # Append the user message to the session state
                    st.session_state.messages.append({"role": "user", "content": prompt})
                    

                    # Append the document text to the context (if available)
                    messages_with_context = st.session_state.messages.copy()
                    
                    messages_with_context.append({"role": "system", "content": f"{editable_prompt}. {context}"})

                    # Call the OpenAI API to generate a response
                    try:
                        client = OpenAI(api_key=openai_api_key)  # Set the API key
                        response = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=messages_with_context
                        )
                        # Extract the message from the response
                        msg = response.choices[0].message.content

                        # Append the assistant's message to the session state
                        st.session_state['messages'].append({"role": "assistant", "content": msg})
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
                    return msg

                response_container = st.container()
                # container for text box
                container = st.container()

                with container:
                    with st.form(key='my_form', clear_on_submit=True):
                        user_input = st.text_area("You:", key='input', height=100)
                        submit_button = st.form_submit_button(label='Send')

                    if submit_button and user_input:
                        output = generate_response(user_input)
                        st.session_state['past'].append(user_input)
                        st.session_state['generated'].append(output)   

                if st.session_state['generated']:
                    with response_container:
                        for i in range(len(st.session_state['generated'])):
                            message(st.session_state["past"][i], is_user=True, key=str(i) + '_user', avatar_style="initials", seed="ME")
                            message(st.session_state["generated"][i], key=str(i),avatar_style="initials", seed="AI")

            with tab3:
                if uploaded_file:
                    st.subheader("Strategy recommendations")
                    st.write(recommendations)

                    st.write(ppt)
                else:
                    st.info("Upload a document to see the summary")

elif st.session_state.selected =="Data cleaning":
        openai_api_key = st.text_input("Put your api key in here and press enter")


        st.title("Creating the shimpnet file")
        bidfile = st.file_uploader("Upload an Excel file", type=["xlsx", "xls"])
        if bidfile is not None:
            bid_df = pd.read_excel(bidfile)
            bid_df.columns = bid_df.columns.str.strip()
            bd_head=bid_df.head()
            st.write(bd_head)
            bd_text= bd_head.to_csv(index=False)
            st.write(bd_text)






                                               # checking if there's is ZC in seperate column #

            # request = [{"role": "system", "content":f"I have a pandas DataFrame named 'bid_df'. Check if it contains any columns related to zip codes or postal codes. If such columns exist, create a variable 'column_zc' and set it to True, then rename the column related to the origin/collection zip code to 'origin ZC' and the column related to the delivery/destination zip code to 'destination ZC'. If no zip code-related columns are found, set 'column_zc' to False. Return only the Python code (no explanations). This is the head of my DataFrame: {bd_text}" }]
            # client = OpenAI(api_key=openai_api_key) 
            # response = client.chat.completions.create(
            #     model="gpt-4o-mini",
            #     temperature= 0.0,
            #     messages=request)
            # code_zc_colum = response.choices[0].message.content
            # st.code(code_zc_colum, language="python")
            # code_zc_colum = code_zc_colum.strip("```python").strip("```").strip()
            # exec(code_zc_colum)  
    
            # if column_zc:
            #     st.write(bid_df) 







                                                                                     # checking if there's is iso #

            request = [{"role": "system", "content":f"I have a pandas DataFrame named 'bid_df'. Do you see any columns that has the country ISO2 abbreviations (e.g., 'FR', 'DE')? If you do, return iso==True . If no such column is found, return iso=False .Please be very very accurate. Return only what's asked (no explanations) . This is the head of my DataFrame: {bd_text}" }]
            client = OpenAI(api_key=openai_api_key) 
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                temperature= 0.0,
                messages=request)
            code_iso = response.choices[0].message.content
            st.code(code_iso, language="python")
            code_iso = code_iso.strip("```python").strip("```").strip()
            exec(code_iso)  
            st.write(bid_df) 


                                                                                       ##### checking if column for origin country exist###
            request = [{"role": "system", "content":f"I have a pandas DataFrame named 'bid_df'. Carefully analyze the column names and their corresponding data to determine if there is a column that represents the **origin country** (not destination, not city, not post code, not country code). The column should contain **country names** (e.g., 'France', 'Spain', 'Germany') and should have a name related to origin, such as 'origin', 'collection', 'depart','chargement', 'country from', 'shipping country','consignor'..... or another similar name that strongly indicates the **origin country**. You must be at least 95% certain that the column represents the origin country and not the destination.If so, return a Python script to rename that column to 'origin_cntry'. Else, return a variable 'origin_country_existing'=False. Do not guess. Do not assume. Do not include explanations.Do not check for the existence of the columns in the code, just return the renaming code directly Return only the Python code (no explanations). This is the head of my DataFrame: {bd_text}" }]
            client = OpenAI(api_key=openai_api_key) 
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                temperature= 0,
                top_p = 0.1,
                messages=request)
            code_rename_country = response.choices[0].message.content
            st.code(code_rename_country, language="python")
            code_rename_country = code_rename_country.strip("```python").strip("```").strip()
            exec(code_rename_country)  
            st.write(bid_df)
                                                                                ##### checking if column for destination country exist###
            request = [{"role": "system", "content":f"I have a pandas DataFrame named 'bid_df'. Carefully analyze the column names and their corresponding data to determine if there is a column that represents the **destination country** (not origin, not city, not post code). The column should contain **country names** (e.g., 'France', 'Spain', 'Germany') and should have a name related to destination, such as 'destination', 'arrivée', 'country to','livraison','consignee'...., 'delivery' or another similar name that strongly indicates the **destination country**. You must be at least 95% certain that the column represents the destination country and not the origin.If so, return a Python script to rename that column to 'destination_cntry'. Else, return a variable 'destination_country_existing'=False. Do not guess. Do not assume. Do not include explanations.Do not check for the existence of the columns in the code, just return the renaming code directly Return only the Python code (no explanations). This is the head of my DataFrame: {bd_text}" }]
            client = OpenAI(api_key=openai_api_key) 
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                temperature= 0,
                top_p = 0.1,
                messages=request)
            code_rename_country = response.choices[0].message.content
            st.code(code_rename_country, language="python")
            code_rename_country = code_rename_country.strip("```python").strip("```").strip()
            exec(code_rename_country)  
            st.write(bid_df)


            if iso==False:
                iso_data=pd.read_excel("ISO2.xlsx")
                melted_iso_data = iso_data.melt(id_vars=["ISO2"], value_name="Country").drop(columns=["variable"])
                melted_iso_data = melted_iso_data.drop_duplicates(subset=["Country"])
                def fuzzy_match(value, choices,score=78):
                    match = process.extractOne(value, choices)  # Returns (matched_value, score, etc.)
                    if match and match[1>= score]:
                        return match[0],match[1]  # Extract only the matched value
                    return None

                if "origin_cntry" in bid_df.columns:
                    bid_df[["matched_country", "origin_country_match_score"]] = bid_df["origin_cntry"].apply(lambda x: pd.Series(fuzzy_match(x, melted_iso_data["Country"])))
                    bid_df = bid_df.merge(melted_iso_data, left_on="matched_country", right_on="Country", how="left")
                    bid_df=bid_df.drop(columns=["Country","matched_country"]).rename(columns={"ISO2": "iso_origin"})
                    st.write(bid_df)

                
                    
                if "destination_cntry" in bid_df.columns:
                    bid_df[["matched_country", "delivery_country_match_score"]] = bid_df["destination_cntry"].apply(lambda x: pd.Series(fuzzy_match(x, melted_iso_data["Country"])))
                    bid_df = bid_df.merge(melted_iso_data, left_on="matched_country", right_on="Country", how="left")
                    bid_df=bid_df.drop(columns=["Country","matched_country"]).rename(columns={"ISO2": "iso_origin"})
                    st.write(bid_df)


                
            # if iso==True:
            #         request = [{"role": "system", "content":f"I have a pandas DataFrame named 'bid_df'. Carefully analyze the column names and their corresponding data to determine if there is a column that represents the  **ISO2** country code (e.g. FR , DE, ES...) for origin county( do not confuse with iso country code of the destination country).The column should contain **country  iso2 code** (e.g., 'FR', 'ES', 'DE') and should have a name related to origin, such as 'origin', 'collection', 'depart','chargement', 'country from', 'shipping country','consignor' or another similar name that strongly indicates the **origin country**. You must be at least 95% certain that the column represents the origin country and not the destination.If so, return a Python script to rename that column to 'origin_cntry'. Else, return a variable 'origin_country_existing'=False. Do not guess. Do not assume. Do not include explanations.Do not check for the existence of the columns in the code, just return the renaming code directly Return only the Python code (no explanations). This is the head of my DataFrame: {bd_text}" }]
            #         client = OpenAI(api_key=openai_api_key) 
            #         response = client.chat.completions.create(
            #             model="gpt-4o-mini",
            #             temperature= 0.0,
            #             messages=request)
            #         code_rename_country = response.choices[0].message.content
            #         st.code(code_rename_country, language="python")
            #         code_rename_country = code_rename_country.strip("```python").strip("```").strip()
            #         exec(code_rename_country)  
            #         st.write(bid_df)



                
                











            # request = [{"role": "system", "content":f"I'll give you the head of a panads dataframe called bid_df, i want you to write me only a python code (i'll excute it directly so just return a code) that create in the same dataframe a new column called 'ZC from' that takes the origin country code (2 letters) and the first two digit of the origin zip or post code. and the same for a column called 'ZC to' (for the destination).here's an example: FR 92 ( keep in mind the df is not clean may have missing value or not str or not integer or numbers and letters at the same time .. Avoid any possible error) . This is my dataframe: {bd_text} " }]
            # try:
            #     client = OpenAI(api_key=openai_api_key) 
            #     response = client.chat.completions.create(
            #         model="gpt-4o-mini",
            #         messages=request)
            #     code = response.choices[0].message.content
            # except Exception as e:
            #     st.error(f"An error occurred: {e}")

            # code = code.strip("```python").strip("```").strip()

            # st.write("Code after cleaning:")
            # st.code(code, language="python")

            # exec(code)  
            # st.write(bid_df)  





